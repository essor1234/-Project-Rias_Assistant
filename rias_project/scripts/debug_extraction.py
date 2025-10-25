# generate_edu_materials.py
import json
import sys
import datetime
import time
from pathlib import Path
from dotenv import load_dotenv
from openai import OpenAI
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from zipfile import ZipFile

# ----------------------------------------------------------------------
# LOGGING SETUP — print to both console and log file
# ----------------------------------------------------------------------
class Tee:
    def __init__(self, *files): self.files = files
    def write(self, obj):
        for f in self.files: f.write(obj); f.flush()
    def flush(self):
        for f in self.files: f.flush()

timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
SCRIPT_DIR = Path(__file__).resolve().parent.parent
LOG_DIR = SCRIPT_DIR / "logs"
LOG_DIR.mkdir(exist_ok=True)
LOG_FILE = LOG_DIR / f"edu_materials_log_{timestamp}.txt"
log_f = open(LOG_FILE, "w", encoding="utf-8")
sys.stdout = Tee(sys.__stdout__, log_f)
sys.stderr = Tee(sys.__stderr__, log_f)
# ----------------------------------------------------------------------

load_dotenv()
client = OpenAI()

# ----------------------------------------------------------------------
# CONFIGURATION
# ----------------------------------------------------------------------
TXT_DIR       = SCRIPT_DIR / "data" / "extracted_text" / "test4"
PROMPT_PATH   = SCRIPT_DIR / "prompts" / "[Prompt]explain_and_lab.txt"
MODEL         = "gpt-5"          # ✅ switched from gpt-4o to gpt-5
MAX_TOKENS    = 20000
TEMPERATURE   = 1
MAX_RETRIES   = 3
# ----------------------------------------------------------------------


def truncate_text(text: str, limit: int = 20000) -> str:
    return text if len(text) <= limit else text[:limit] + "\n\n[Text truncated for LLM]"


def load_prompt() -> str:
    return PROMPT_PATH.read_text(encoding="utf-8")


# ----------------------------------------------------------------------
# LLM CALL — updated to GPT-5 format
# ----------------------------------------------------------------------
def call_llm(prompt_text: str) -> str:
    """
    Call GPT-5 with robust retry and enforced JSON schema.
    Output schema:
    {
      "Slides": [...],
      "Labs": [...]
    }
    """
    system_msg = {
        "role": "system",
        "content": (
            "You are an educational AI tutor that outputs structured JSON only.\n"
            "Output schema:\n"
            "{ \"Slides\": [ {\"Title\": str, \"Content\": str, ...} ], "
            "\"Labs\": [ {\"Title\": str, \"Dataset\": {...}, \"CodeFiles\": [...] } ] }\n"
            "Never output markdown or commentary — JSON only."
        ),
    }

    user_msg = {"role": "user", "content": prompt_text}

    for attempt in range(MAX_RETRIES):
        try:
            resp = client.chat.completions.create(
                model=MODEL,
                messages=[system_msg, user_msg],
                max_completion_tokens=MAX_TOKENS,
                # temperature=TEMPERATURE,
                response_format={"type": "json_object"},
            )

            print(f"LLM response (attempt {attempt+1}) OK")
            choice = resp.choices[0]
            content = getattr(choice.message, "content", None)
            if not content:
                print("⚠️ Empty content from GPT-5")
                continue

            cleaned = clean_raw(content)
            return cleaned.strip()

        except Exception as e:
            print(f"Attempt {attempt+1}/{MAX_RETRIES} failed: {e}")
            if attempt == MAX_RETRIES - 1:
                raise
            time.sleep(2 ** attempt)
    return ""


def clean_raw(raw: str) -> str:
    raw = raw.strip()
    if raw.startswith("```"):
        raw = raw.split("```", 1)[1].rsplit("```", 1)[0]
    if raw.lower().startswith("json"):
        raw = raw[4:].lstrip()
    return raw.strip()


# ----------------------------------------------------------------------
# PowerPoint Writer
# ----------------------------------------------------------------------
def create_ppt(slides, output_path, raw_json_sample=None):
    prs = Presentation()
    for slide in slides:
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        s = prs.slides.add_slide(layout)

        # Title
        title = slide.get("Title") or slide.get("Heading") or "Untitled"
        try:
            s.shapes.title.text = str(title)
        except Exception:
            pass

        # Content
        if slide.get("Content"):
            content_text = slide.get("Content")
        else:
            parts = []
            for k in ("Equation", "ConceptExplanation", "DeepExplanation", "RealExample", "ImageIdea"):
                v = slide.get(k)
                if v: parts.append(f"{k}: {v}")
            content_text = "\n\n".join(parts) if parts else ""

        # Text placeholder
        body_shape = None
        for shape in s.shapes:
            if hasattr(shape, "text_frame") and shape != s.shapes.title:
                body_shape = shape
                break

        if body_shape is None:
            left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)
            body_shape = s.shapes.add_textbox(left, top, width, height)

        body_shape.text_frame.clear()
        p = body_shape.text_frame.paragraphs[0]
        p.text = content_text[:4000] or "No detailed content available."

    prs.save(output_path)
    print(f"📊 Slides saved → {output_path}")

    if raw_json_sample:
        raw_path = Path(output_path).with_suffix(".raw.json")
        raw_path.write_text(json.dumps(raw_json_sample, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"🔎 Raw JSON saved → {raw_path}")


# ----------------------------------------------------------------------
# ZIP Generator for Labs
# ----------------------------------------------------------------------
def create_lab_zip(labs, output_zip_path):
    if not labs:
        print("⚠️ No labs to package.")
        return

    with ZipFile(output_zip_path, "w") as zf:
        for i, lab in enumerate(labs, start=1):
            title = lab.get("Title", f"Lab_{i}")
            dataset = lab.get("Dataset", {})
            codefiles = lab.get("CodeFiles", [])

            if dataset:
                csv_name = dataset.get("filename", f"{title.replace(' ', '_')}.csv")
                csv_content = "x,y,true_label,pred_label\n1,0.8,1,1\n2,0.3,1,0\n3,0.9,0,1\n"
                zf.writestr(csv_name, csv_content)
                readme = f"# Dataset: {csv_name}\n\n{dataset.get('description', '')}\n"
                zf.writestr(f"{title}_README.txt", readme)

            for j, cfile in enumerate(codefiles, start=1):
                name = cfile.get("filename", f"exercise_{i}_{j}.py")
                desc = cfile.get("description", "")
                code_content = (
                    f"# {desc}\n"
                    f"import pandas as pd\nimport numpy as np\nimport matplotlib.pyplot as plt\n\n"
                    f"data = pd.read_csv('{dataset.get('filename','data.csv')}')\n"
                    f"print(data.head())\n"
                )
                zf.writestr(name, code_content)

    print(f"🧩 Lab files saved → {output_zip_path}")


# ----------------------------------------------------------------------
# Main Processing
# ----------------------------------------------------------------------
def process_paper(txt_path, base_prompt, edu_output):
    print(f"\n📘 Processing {txt_path.name}")
    text = truncate_text(txt_path.read_text(encoding="utf-8"))
    combined = base_prompt.replace("<<<DOCUMENT_TEXT>>>", text)

    raw = call_llm(combined)
    cleaned = clean_raw(raw)

    if not cleaned:
        print("⚠️ No response from GPT-5.")
        return [], []

    debug_raw_path = edu_output / f"{txt_path.stem}_raw.txt"
    debug_raw_path.write_text(cleaned, encoding="utf-8")

    try:
        data = json.loads(cleaned)
        slides = data.get("Slides", [])
        labs = data.get("Labs", [])
        print(f"✅ Parsed JSON: {len(slides)} slides, {len(labs)} labs")
        return slides, labs, data
    except json.JSONDecodeError as e:
        print(f"❌ JSON parse error: {e}")
        print("Raw output snippet:", cleaned[:1000])
        return [], [], {}


def main():
    print("\n--- Generating Educational Materials (GPT-5) ---")
    prompt = load_prompt()
    txt_files = sorted(TXT_DIR.glob("*.txt"))
    if not txt_files:
        print("❌ No .txt files found")
        return

    edu_output = SCRIPT_DIR / "data" / "edu_output"
    edu_output.mkdir(parents=True, exist_ok=True)

    for txt_path in txt_files:
        slides, labs, data = process_paper(txt_path, prompt, edu_output)
        if not slides and not labs:
            print(f"⚠️ No educational content generated for {txt_path.name}")
            continue

        out_ppt = edu_output / f"slides_{txt_path.stem}.pptx"
        out_zip = edu_output / f"lab_{txt_path.stem}.zip"

        if slides:
            create_ppt(slides, out_ppt, raw_json_sample=data)
        if labs:
            create_lab_zip(labs, out_zip)

    print("\n✅ All papers processed successfully!")
    print(f"Log file saved → {LOG_FILE}")


if __name__ == "__main__":
    main()
