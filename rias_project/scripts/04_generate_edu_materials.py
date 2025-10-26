# edu_materials_generator.py
import json
import sys
import datetime
from pathlib import Path
from dotenv import load_dotenv
from openai import OpenAI
import time
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from zipfile import ZipFile
from typing import List, Dict, Any, Tuple, Optional


# ----------------------------------------------------------------------
# LOGGING SETUP — print to both console and log file
# ----------------------------------------------------------------------
class Tee:
    def __init__(self, *files):
        self.files = files

    def write(self, obj):
        for f in self.files:
            f.write(obj)
            f.flush()

    def flush(self):
        for f in self.files:
            f.flush()


class EducationalMaterialsGenerator:
    def __init__(
        self,
        script_dir: Path = None,
        txt_dir: str = "data/extracted_text/test4",
        prompt_path: str = "prompts/[Prompt]explain_and_lab.txt",
        output_dir: str = "data/edu_output",
        model: str = "gpt-4o",
        max_tokens: int = 10000,
        temperature: float = 0.0,
        max_retries: int = 3,
        text_limit: int = 20_000,
    ):
        # Setup paths
        self.SCRIPT_DIR = script_dir or Path(__file__).resolve().parent.parent
        self.TXT_DIR = self.SCRIPT_DIR / txt_dir
        self.PROMPT_PATH = self.SCRIPT_DIR / prompt_path
        self.EDU_OUTPUT = self.SCRIPT_DIR / output_dir
        self.EDU_OUTPUT.mkdir(parents=True, exist_ok=True)

        # Config
        self.MODEL = model
        self.MAX_TOKENS = max_tokens
        self.TEMPERATURE = temperature
        self.MAX_RETRIES = max_retries
        self.TEXT_LIMIT = text_limit

        # Logging
        self._setup_logging()

        # OpenAI
        load_dotenv()
        self.client = OpenAI()

        # Load prompt once
        self.base_prompt = self._load_prompt()

        print(f"EduMaterialsGenerator initialized.")
        print(f"Input: {self.TXT_DIR}")
        print(f"Output: {self.EDU_OUTPUT}")

    def _setup_logging(self):
        timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        log_dir = self.SCRIPT_DIR / "logs"
        log_dir.mkdir(exist_ok=True)
        log_file = log_dir / f"edu_materials_log_{timestamp}.txt"
        log_f = open(log_file, "w", encoding="utf-8")
        sys.stdout = Tee(sys.__stdout__, log_f)
        sys.stderr = Tee(sys.__stderr__, log_f)
        self.log_file = log_file
        print(f"Logging to {self.log_file}")

    def _load_prompt(self) -> str:
        if not self.PROMPT_PATH.exists():
            raise FileNotFoundError(f"Prompt file not found: {self.PROMPT_PATH}")
        return self.PROMPT_PATH.read_text(encoding="utf-8")

    def _truncate_text(self, text: str) -> str:
        return text if len(text) <= self.TEXT_LIMIT else text[: self.TEXT_LIMIT] + "\n\n[Text truncated for LLM]"

    def _call_llm(self, prompt_text: str) -> str:
        for attempt in range(self.MAX_RETRIES):
            try:
                resp = self.client.chat.completions.create(
                    model=self.MODEL,
                    messages=[
                        {
                            "role": "system",
                            "content": (
                                "You are an educational AI tutor that outputs structured JSON only. "
                                "If the paper has few metrics, you may generate example metrics for teaching, "
                                "but keep consistent JSON keys. "
                                "Never output markdown or commentary — JSON only."
                            ),
                        },
                        {"role": "user", "content": prompt_text},
                    ],
                    max_tokens=self.MAX_TOKENS,
                    temperature=self.TEMPERATURE,
                    response_format={"type": "json_object"},
                )
                return resp.choices[0].message.content.strip()
            except Exception as e:
                print(f"Attempt {attempt + 1} failed: {e}")
                if attempt == self.MAX_RETRIES - 1:
                    raise
                time.sleep(2 ** attempt)
        return ""

    def _clean_raw(self, raw: str) -> str:
        if not raw:
            return ""
        if raw.startswith("```"):
            parts = raw.split("```", 2)
            raw = parts[1] if len(parts) > 2 else parts[0]
        raw = raw.strip()
        if raw.lower().startswith("json"):
            raw = raw[4:].strip()
        return raw

    def _create_ppt(self, slides: List[Dict], output_path: Path, raw_json_sample: Dict = None):
        prs = Presentation()
        for slide in slides:
            layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            s = prs.slides.add_slide(layout)

            # Title
            title = slide.get("Title") or slide.get("Heading") or "Untitled"
            try:
                s.shapes.title.text = str(title)
            except Exception:
                pass

            # Content
            if slide.get("Content"):
                content_text = slide.get("Content")
            else:
                parts = []
                for k in ("Equation", "ConceptExplanation", "DeepExplanation", "RealExample", "ImageIdea"):
                    v = slide.get(k)
                    if v:
                        parts.append(f"{k}: {v}")
                content_text = "\n\n".join(parts) if parts else ""

            # Find body placeholder
            body_shape = None
            for shape in s.shapes:
                if hasattr(shape, "text_frame") and shape != s.shapes.title:
                    body_shape = shape
                    break

            if body_shape is None:
                left = top = Inches(0.5)
                width = Inches(9)
                height = Inches(4.5)
                body_shape = s.shapes.add_textbox(left, top, width, height)

            body_shape.text_frame.clear()
            content_text = content_text[:4000] or "No detailed content available."
            p = body_shape.text_frame.paragraphs[0]
            p.text = content_text

        prs.save(output_path)
        print(f"Slides saved to {output_path}")

        if raw_json_sample:
            raw_path = output_path.with_suffix(".raw.json")
            raw_path.write_text(json.dumps(raw_json_sample, ensure_ascii=False, indent=2), encoding="utf-8")
            print(f"Raw JSON saved to {raw_path}")

    def _create_lab_zip(self, labs: List[Dict], output_zip_path: Path):
        if not labs:
            print("No labs to package.")
            return

        with ZipFile(output_zip_path, "w") as zf:
            for i, lab in enumerate(labs, start=1):
                title = lab.get("Title", f"Lab_{i}")
                dataset = lab.get("Dataset", {})
                codefiles = lab.get("CodeFiles", [])

                # CSV
                if dataset:
                    csv_name = dataset.get("filename", f"{title.replace(' ', '_')}.csv")
                    csv_content = "x,y,true_label,pred_label\n1,0.8,1,1\n2,0.3,1,0\n3,0.9,0,1\n"
                    zf.writestr(csv_name, csv_content)
                    readme = f"# Dataset: {csv_name}\n\n{dataset.get('description', '')}\n"
                    zf.writestr(f"{title}_README.txt", readme)

                # Python files
                for j, cfile in enumerate(codefiles, start=1):
                    name = cfile.get("filename", f"exercise_{i}_{j}.py")
                    desc = cfile.get("description", "")
                    code_content = (
                        f"# {desc}\n"
                        f"import pandas as pd\nimport numpy as np\nimport matplotlib.pyplot as plt\n\n"
                        f"data = pd.read_csv('{dataset.get('filename','data.csv')}')\n"
                        f"print(data.head())\n"
                    )
                    zf.writestr(name, code_content)

        print(f"Lab files saved to {output_zip_path}")

    def process_single(self, txt_path: Path) -> Tuple[List[Dict], List[Dict], Dict]:
        """Process one .txt file and return slides, labs, raw data."""
        print(f"\nProcessing {txt_path.name}")
        text = self._truncate_text(txt_path.read_text(encoding="utf-8"))
        prompt = self.base_prompt.replace("<<<DOCUMENT_TEXT>>>", text)

        raw = self._call_llm(prompt)
        cleaned = self._clean_raw(raw)

        if not cleaned:
            print("No response from GPT.")
            return [], [], {}

        debug_path = self.EDU_OUTPUT / f"{txt_path.stem}_raw.txt"
        debug_path.write_text(cleaned, encoding="utf-8")

        try:
            data = json.loads(cleaned)
            slides = data.get("Slides", [])
            labs = data.get("Labs", [])
            print(f"Parsed: {len(slides)} slides, {len(labs)} labs")
            return slides, labs, data
        except json.JSONDecodeError as e:
            print(f"JSON parse error: {e}")
            print("Raw:", cleaned[:1000])
            return [], [], {}

    def generate_all(self) -> List[Dict[str, Any]]:
        """Process all .txt files in TXT_DIR and generate PPTX + ZIP."""
        txt_files = sorted(self.TXT_DIR.glob("*.txt"))
        if not txt_files:
            print("No .txt files found")
            return []

        results = []
        for txt_path in txt_files:
            slides, labs, data = self.process_single(txt_path)

            if not slides and not labs:
                print(f"No content generated for {txt_path.name}")
                continue

            out_ppt = self.EDU_OUTPUT / f"slides_{txt_path.stem}.pptx"
            out_zip = self.EDU_OUTPUT / f"lab_{txt_path.stem}.zip"

            if slides:
                self._create_ppt(slides, out_ppt, raw_json_sample=data)
            if labs:
                self._create_lab_zip(labs, out_zip)

            results.append({
                "file": txt_path.name,
                "pptx": out_ppt.name if slides else None,
                "zip": out_zip.name if labs else None,
                "slides_count": len(slides),
                "labs_count": len(labs),
            })

        print("\nAll papers processed successfully!")
        print(f"Log file: {self.log_file}")
        return results
    

# ...existing code...

# This function should REPLACE the old `def run(...)` 
# at the end of your 'scripts/01_extract_text.py' file.

def run(pdf_path, out_dir, prev=None):
    """Bridge function for main.py pipeline."""
    try:
        p = Path(pdf_path)
        out = Path(out_dir)
        
        # Get project root for paths
        SCRIPT_DIR = Path(__file__).resolve().parent.parent
        
        # Create generator with script_dir
        generator = EducationalMaterialsGenerator(
            script_dir=SCRIPT_DIR,
            txt_dir=f"data/extracted_text/{p.stem}",
            prompt_path="prompts/[Prompt]explain_and_lab.txt",
            output_dir=str(out)
        )
        
        # Generate materials
        results = generator.generate_all()
        
        # Return files created
        files = [f.name for f in out.glob("*") if f.is_file()]
        return {
            "status": "success",
            "files": files,
            "summary": "materials generated"
        }
        
    except Exception as e:
        print(f"ERROR in generate_edu: {e}")
        return {"status": "error", "error": str(e)}

